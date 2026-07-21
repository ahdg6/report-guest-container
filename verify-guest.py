from __future__ import annotations

import importlib.metadata
import json
import os
from pathlib import Path
import shutil
import socket
import subprocess

import fitz
import pandas
import pyxlsb
from docx import Document
from odf import teletype
from odf.opendocument import OpenDocumentSpreadsheet, load
from odf.table import Table, TableCell, TableRow
from odf.text import P
from openpyxl import Workbook, load_workbook
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation


ROOT = Path.cwd()
MARKER = "PLUXEL_GUEST_OK_7319"
CAPABILITIES_PATH = Path("/opt/pluxel/guest-capabilities.json")


def run(*argv: str) -> str:
    completed = subprocess.run(
        argv,
        cwd=ROOT,
        check=True,
        capture_output=True,
        text=True,
        timeout=30,
    )
    return completed.stdout


def assert_runtime_boundary() -> None:
    assert os.geteuid() != 0, "guest commands must not run as root"
    try:
        with socket.create_connection(("1.1.1.1", 80), timeout=1):
            raise AssertionError("guest network policy unexpectedly allowed an outbound connection")
    except OSError:
        pass


def verify_capabilities() -> dict[str, object]:
    capabilities = json.loads(CAPABILITIES_PATH.read_text(encoding="utf-8"))
    assert capabilities["schemaVersion"] == 1
    for command in capabilities["commands"]:
        assert shutil.which(command), f"missing guest command: {command}"
    for module in capabilities["pythonModules"]:
        __import__(module)
    return capabilities


def verify_navigation() -> None:
    (ROOT / "records.json").write_text(
        json.dumps({"marker": MARKER, "records": [{"amount": 73}, {"amount": 19}]}),
        encoding="utf-8",
    )
    assert run("jq", "-r", ".marker", "records.json").strip() == MARKER
    assert MARKER in run("rg", "--fixed-strings", MARKER, "records.json")
    assert "JSON" in run("file", "--brief", "records.json")


def verify_workspace_versioning(capabilities: dict[str, object]) -> None:
    versioning = capabilities["workspaceVersioning"]
    assert isinstance(versioning, dict)
    assert versioning["command"] == "jj"
    assert versioning["version"] == "0.43.0"
    assert versioning["backend"] == "git"
    assert versioning["remoteRequired"] is False
    assert run("jj", "--version").startswith("jj 0.43.0")
    guide = Path(versioning["agentGuide"]).read_text(encoding="utf-8")
    assert "jj describe" in guide
    assert "jj new" in guide
    assert "不要自行执行 `jj undo`" in guide

    workspace = ROOT / "jj-workspace"
    protected_directories = ("sources", "evidence", "searches", "templates", "tools")
    for directory in protected_directories:
        protected = workspace / directory
        protected.mkdir(parents=True)
        (protected / "ignored.md").write_text("read-only projection\n", encoding="utf-8")
    internal = workspace / ".report-agent"
    internal.mkdir(parents=True)
    (internal / "state.json").write_text("{}\n", encoding="utf-8")
    (workspace / ".report-agent-state.json").write_text("{}\n", encoding="utf-8")
    (workspace / "existing.md").write_text("# Existing workspace\n", encoding="utf-8")
    (workspace / "workbook.xlsx").write_bytes(b"PK\x03\x04writable workbook checkpoint")
    run("jj", "git", "init", "--no-colocate", str(workspace))
    assert (workspace / ".jj").is_dir()
    assert not (workspace / ".git").exists()
    initial_status = run("jj", "-R", str(workspace), "status")
    assert "existing.md" in initial_status
    assert "workbook.xlsx" in initial_status
    for ignored in (*protected_directories, ".report-agent"):
        assert f"{ignored}/" not in initial_status
    assert ".report-agent-state.json" not in initial_status
    tracked = run("jj", "-R", str(workspace), "file", "list")
    assert "existing.md" in tracked
    assert "workbook.xlsx" in tracked
    assert "ignored.md" not in tracked
    run("jj", "-R", str(workspace), "describe", "-m", "Initial workspace")
    run("jj", "-R", str(workspace), "new")
    draft = workspace / "report.md"
    draft.write_text("# Report\n\nFirst draft.\n", encoding="utf-8")
    status = run("jj", "-R", str(workspace), "status")
    assert "report.md" in status
    assert "sources/source.md" not in status
    assert "report.md" in run("jj", "-R", str(workspace), "diff", "--stat")
    run("jj", "-R", str(workspace), "describe", "-m", "first draft")
    run("jj", "-R", str(workspace), "new")
    draft.write_text("# Report\n\nRevised draft.\n", encoding="utf-8")
    assert "Revised draft." in run("jj", "-R", str(workspace), "diff")
    assert "first draft" in run("jj", "-R", str(workspace), "log", "--no-graph")
    assert "First draft." in run("jj", "-R", str(workspace), "show", "@-")
    run("jj", "-R", str(workspace), "describe", "-m", "revised draft")
    run("jj", "-R", str(workspace), "new")
    draft.write_text("# Report\n\nAccidental overwrite.\n", encoding="utf-8")
    run("jj", "-R", str(workspace), "restore", "--from", "@-", "root-file:report.md")
    assert "Revised draft." in draft.read_text(encoding="utf-8")


def verify_pdf() -> None:
    path = ROOT / "sample.pdf"
    document = fitz.open()
    page = document.new_page()
    page.insert_text((72, 72), MARKER)
    document.save(path)
    document.close()
    assert "Pages:           1" in run("pdfinfo", path.name)
    run("pdftotext", path.name, "sample.txt")
    assert MARKER in (ROOT / "sample.txt").read_text(encoding="utf-8")
    reopened = fitz.open(path)
    assert MARKER in "".join(page.get_text() for page in reopened)
    assert reopened[0].get_pixmap().width > 0
    reopened.close()


def verify_word_and_pandoc() -> None:
    path = ROOT / "sample.docx"
    document = Document()
    document.add_heading("Pluxel guest", level=1)
    document.add_paragraph(MARKER)
    table = document.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "amount"
    table.cell(0, 1).text = "7319"
    document.save(path)
    reopened = Document(path)
    assert MARKER in "\n".join(paragraph.text for paragraph in reopened.paragraphs)
    assert MARKER in run("pandoc", path.name, "--to=plain")


def verify_powerpoint() -> None:
    path = ROOT / "sample.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = MARKER
    presentation.save(path)
    reopened = Presentation(path)
    text = "\n".join(
        shape.text
        for current_slide in reopened.slides
        for shape in current_slide.shapes
        if hasattr(shape, "text")
    )
    assert MARKER in text


def verify_spreadsheets() -> None:
    xlsx = ROOT / "sample.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Checks"
    sheet.append(["marker", "amount", "double"])
    sheet.append([MARKER, 7319, "=B2*2"])
    workbook.save(xlsx)
    formulas = load_workbook(xlsx, data_only=False)
    assert formulas["Checks"]["C2"].value == "=B2*2"
    frame = pandas.read_excel(xlsx, sheet_name="Checks")
    assert frame.iloc[0]["marker"] == MARKER

    ods = ROOT / "sample.ods"
    spreadsheet = OpenDocumentSpreadsheet()
    table = Table(name="Checks")
    row = TableRow()
    cell = TableCell(valuetype="string")
    cell.addElement(P(text=MARKER))
    row.addElement(cell)
    table.addElement(row)
    spreadsheet.spreadsheet.addElement(table)
    spreadsheet.save(str(ods), addsuffix=False)
    reopened = load(str(ods))
    assert MARKER in teletype.extractText(reopened.spreadsheet)
    assert importlib.metadata.version("pyxlsb") == pyxlsb.__version__


def verify_image_and_ocr() -> None:
    path = ROOT / "sample.png"
    image = Image.new("RGB", (1100, 220), "white")
    font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 72)
    ImageDraw.Draw(image).text((35, 60), "PLUXEL OCR 7319", fill="black", font=font)
    image.save(path)
    reopened = Image.open(path)
    assert reopened.size == (1100, 220)
    ocr = run("tesseract", path.name, "stdout", "--psm", "7", "-l", "eng")
    assert "PLUXEL" in ocr.upper() and "7319" in ocr


def verify_archives() -> None:
    marker = ROOT / "archive-marker.txt"
    marker.write_text(MARKER, encoding="utf-8")
    run("zip", "-q", "sample.zip", marker.name)
    assert run("unzip", "-p", "sample.zip", marker.name).strip() == MARKER
    run("7z", "a", "-bd", "-y", "sample.7z", marker.name)
    assert MARKER in run("7z", "e", "-so", "sample.7z", marker.name)


def main() -> None:
    assert_runtime_boundary()
    capabilities = verify_capabilities()
    verify_navigation()
    verify_workspace_versioning(capabilities)
    verify_pdf()
    verify_word_and_pandoc()
    verify_powerpoint()
    verify_spreadsheets()
    verify_image_and_ocr()
    verify_archives()
    result = {
        "status": "ok",
        "marker": MARKER,
        "capabilitySchemaVersion": capabilities["schemaVersion"],
        "python": os.sys.version.split()[0],
    }
    (ROOT / "guest-smoke-result.json").write_text(
        json.dumps(result, sort_keys=True) + "\n", encoding="utf-8"
    )
    print(json.dumps(result, sort_keys=True))


if __name__ == "__main__":
    main()
